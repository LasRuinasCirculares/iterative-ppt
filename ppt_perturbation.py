import random
import copy
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE


class PPTPerturbation:
    def __init__(self, ppt_path):
        """
        初始化PPT扰动工具
        
        Args:
            ppt_path: PPT文件路径
        """
        self.ppt_path = ppt_path
        self.prs = Presentation(ppt_path)
        # 幻灯片尺寸（EMU单位）
        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height
    
    def delete_random_elements(self, delete_ratio=0.2, preserve_title=True):
        """
        随机删除PPT中的元素（全局删除，不是按幻灯片删除）
        
        Args:
            delete_ratio: 删除比例，默认0.2（20%）
            preserve_title: 是否保留标题，默认True
        
        Returns:
            删除的元素数量
        """
        deleted_count = 0
        
        # 收集所有可以删除的shape（全局）
        all_deletable_shapes = []
        
        for slide in self.prs.slides:
            for shape in slide.shapes:
                # 如果preserve_title为True，只跳过标题占位符
                if preserve_title and shape.is_placeholder:
                    try:
                        placeholder = shape.placeholder_format
                        # 跳过标题类型的占位符
                        if placeholder.type in [1, 3]:  # TITLE = 1, CENTER_TITLE = 3
                            continue
                    except:
                        pass
                
                all_deletable_shapes.append(shape)
        
        # 计算要删除的数量（全局）
        num_to_delete = int(len(all_deletable_shapes) * delete_ratio)
        
        if num_to_delete > 0:
            # 随机选择要删除的shape
            shapes_to_delete = random.sample(all_deletable_shapes, num_to_delete)
            
            # 删除选中的shape
            for shape in shapes_to_delete:
                try:
                    sp = shape.element
                    sp.getparent().remove(sp)
                    deleted_count += 1
                except Exception as e:
                    print(f"删除元素时出错: {e}")
        
        return deleted_count
    
    def create_element_overlaps(self, overlap_ratio=0.3, overlap_intensity=0.5):
        """
        创建元素重叠：移动元素使它们相互重叠（优化版：全局收集元素）
        
        Args:
            overlap_ratio: 要创建重叠的元素比例，默认0.3（30%）
            overlap_intensity: 重叠强度，0-1之间，越大重叠越严重
        
        Returns:
            创建重叠的元素数量
        """
        overlap_count = 0
        
        # 收集所有可移动的shape及其所属的幻灯片
        all_movable_shapes = []
        
        for slide in self.prs.slides:
            for shape in slide.shapes:
                try:
                    if shape.is_placeholder:
                        placeholder = shape.placeholder_format
                        if placeholder.type in [1, 3]:
                            continue
                    if hasattr(shape, 'left') and hasattr(shape, 'top'):
                        all_movable_shapes.append((slide, shape))
                except:
                    if hasattr(shape, 'left') and hasattr(shape, 'top'):
                        all_movable_shapes.append((slide, shape))
        
        if len(all_movable_shapes) < 2:
            return 0
        
        # 计算要创建重叠的元素数量
        num_to_overlap = max(2, int(len(all_movable_shapes) * overlap_ratio))
        
        # 按幻灯片分组
        slide_shapes = {}
        for slide, shape in all_movable_shapes:
            slide_id = id(slide)
            if slide_id not in slide_shapes:
                slide_shapes[slide_id] = []
            slide_shapes[slide_id].append(shape)
        
        # 找出有多个shape的幻灯片
        slides_with_multiple = [shapes for shapes in slide_shapes.values() if len(shapes) >= 2]
        
        if slides_with_multiple:
            # 在有多个shape的幻灯片中创建重叠
            for shapes in slides_with_multiple:
                if overlap_count >= num_to_overlap:
                    break
                
                # 随机选择一些shape进行重叠
                num_to_move = max(2, min(len(shapes), int(len(shapes) * overlap_intensity) + 1))
                shapes_to_overlap = random.sample(shapes, num_to_move)
                
                # 选择一个目标位置（使用第一个shape的位置）
                target_shape = shapes_to_overlap[0]
                target_left = target_shape.left
                target_top = target_shape.top
                
                # 将其他shape移动到附近，造成重叠
                for shape in shapes_to_overlap[1:]:
                    try:
                        # 在目标位置附近随机偏移
                        offset_range = int(min(shape.width, shape.height) * 0.5)
                        offset_x = random.randint(-offset_range, offset_range)
                        offset_y = random.randint(-offset_range, offset_range)
                        
                        shape.left = target_left + offset_x
                        shape.top = target_top + offset_y
                        overlap_count += 1
                    except Exception as e:
                        continue
        else:
            # 如果没有幻灯片有多个shape，就随机移动一些元素到相近位置
            shapes_to_move = random.sample(all_movable_shapes, min(num_to_overlap, len(all_movable_shapes)))
            
            if len(shapes_to_move) >= 2:
                # 选择第一个shape的位置作为目标
                target_slide, target_shape = shapes_to_move[0]
                target_left = target_shape.left
                target_top = target_shape.top
                
                # 将同一幻灯片的其他shape移到附近
                target_slide_id = id(target_slide)
                for slide, shape in shapes_to_move[1:]:
                    if id(slide) == target_slide_id:  # 只移动同一幻灯片的shape
                        try:
                            offset_range = int(min(shape.width, shape.height) * 0.5)
                            offset_x = random.randint(-offset_range, offset_range)
                            offset_y = random.randint(-offset_range, offset_range)
                            
                            shape.left = target_left + offset_x
                            shape.top = target_top + offset_y
                            overlap_count += 1
                        except:
                            continue
        
        return overlap_count
    
    def randomize_element_positions(self, randomize_ratio=0.4, max_shift_ratio=0.3):
        """
        随机移动元素位置，使布局不和谐
        
        Args:
            randomize_ratio: 要随机移动的元素比例，默认0.4（40%）
            max_shift_ratio: 最大位移比例（相对于幻灯片尺寸），默认0.3
        
        Returns:
            移动的元素数量
        """
        moved_count = 0
        
        # 收集所有可移动的shape
        all_movable_shapes = []
        
        for slide in self.prs.slides:
            for shape in slide.shapes:
                try:
                    # 跳过标题
                    if shape.is_placeholder:
                        placeholder = shape.placeholder_format
                        if placeholder.type in [1, 3]:
                            continue
                    
                    if hasattr(shape, 'left') and hasattr(shape, 'top'):
                        all_movable_shapes.append(shape)
                except:
                    if hasattr(shape, 'left') and hasattr(shape, 'top'):
                        all_movable_shapes.append(shape)
        
        # 计算要移动的数量
        num_to_move = int(len(all_movable_shapes) * randomize_ratio)
        
        if num_to_move > 0:
            shapes_to_move = random.sample(all_movable_shapes, num_to_move)
            
            # 计算最大位移量
            max_shift_x = int(self.slide_width * max_shift_ratio)
            max_shift_y = int(self.slide_height * max_shift_ratio)
            
            for shape in shapes_to_move:
                try:
                    # 随机位移
                    shift_x = random.randint(-max_shift_x, max_shift_x)
                    shift_y = random.randint(-max_shift_y, max_shift_y)
                    
                    # 不限制边界，允许部分移出幻灯片（造成布局混乱）
                    new_left = shape.left + shift_x
                    new_top = shape.top + shift_y
                    
                    # 但至少保留一部分在幻灯片内
                    new_left = max(-shape.width // 2, min(new_left, self.slide_width - shape.width // 2))
                    new_top = max(-shape.height // 2, min(new_top, self.slide_height - shape.height // 2))
                    
                    shape.left = new_left
                    shape.top = new_top
                    moved_count += 1
                except Exception as e:
                    continue
        
        return moved_count
    
    def resize_text_boxes(self, resize_ratio=0.3, size_change_range=(0.4, 2.5)):
        """
        调整文本框大小，造成文字溢出或留白过多
        
        Args:
            resize_ratio: 要调整大小的文本框比例，默认0.3（30%）
            size_change_range: 大小变化范围（倍数），默认(0.4, 2.5)
                              小于1会缩小（造成溢出），大于1会放大（造成留白）
        
        Returns:
            调整的文本框数量
        """
        resized_count = 0
        
        # 收集所有文本框
        all_text_shapes = []
        
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    all_text_shapes.append(shape)
        
        # 计算要调整的数量
        num_to_resize = int(len(all_text_shapes) * resize_ratio)
        
        if num_to_resize > 0:
            shapes_to_resize = random.sample(all_text_shapes, num_to_resize)
            
            for shape in shapes_to_resize:
                try:
                    if hasattr(shape, 'width') and hasattr(shape, 'height'):
                        # 随机选择放大或缩小
                        scale_factor = random.uniform(size_change_range[0], size_change_range[1])
                        
                        # 调整宽度和高度
                        original_width = shape.width
                        original_height = shape.height
                        
                        # 随机决定是只改宽度、只改高度，还是同时改
                        change_type = random.choice(['width', 'height', 'both'])
                        
                        if change_type == 'width' or change_type == 'both':
                            new_width = int(original_width * scale_factor)
                            # 确保不会太小或太大
                            new_width = max(int(original_width * 0.3), min(new_width, int(self.slide_width * 0.8)))
                            shape.width = new_width
                        
                        if change_type == 'height' or change_type == 'both':
                            new_height = int(original_height * scale_factor)
                            # 确保不会太小或太大
                            new_height = max(int(original_height * 0.3), min(new_height, int(self.slide_height * 0.8)))
                            shape.height = new_height
                        
                        resized_count += 1
                except Exception as e:
                    print(f"调整文本框大小时出错: {e}")
        
        return resized_count
    
    def change_font_sizes(self, change_ratio=0.3, size_change_range=(-6, 8)):
        """
        随机改变PPT中的字体大小（增强版：更大范围）
        
        Args:
            change_ratio: 改变字体的文本框比例，默认0.3（30%）
            size_change_range: 字体大小改变范围（单位：pt），默认(-6, 8)
        
        Returns:
            改变的文本框数量
        """
        changed_count = 0
        run_changed_count = 0
        
        for slide in self.prs.slides:
            # 收集所有包含文本的shape
            text_shapes = []
            
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():  # 确保有实际文本
                    text_shapes.append(shape)
            
            # 计算要改变字体的shape数量
            num_to_change = max(1, int(len(text_shapes) * change_ratio))  # 至少改变1个
            num_to_change = min(num_to_change, len(text_shapes))  # 不超过总数
            
            if num_to_change > 0 and len(text_shapes) > 0:
                # 随机选择要改变字体大小的shape
                shapes_to_change = random.sample(text_shapes, num_to_change)
                
                for shape in shapes_to_change:
                    shape_changed = False
                    try:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text.strip():  # 确保run有实际文本
                                    # 获取或设置默认字体大小
                                    if run.font.size is not None:
                                        current_size = run.font.size.pt
                                    else:
                                        current_size = 18  # 默认大小
                                    
                                    # 随机改变字体大小
                                    size_change = random.randint(size_change_range[0], size_change_range[1])
                                    if size_change != 0:  # 确保有实际改变
                                        new_size = max(6, min(72, current_size + size_change))  # 6-72pt范围
                                        run.font.size = Pt(new_size)
                                        run_changed_count += 1
                                        shape_changed = True
                        
                        if shape_changed:
                            changed_count += 1
                    except Exception as e:
                        print(f"改变字体大小时出错: {e}")
        
        print(f"   共改变了 {run_changed_count} 个文本运行（run）的字体")
        return changed_count
    
    def apply_combined_perturbations_per_slide(self, delete_ratio=0.2, 
                                               overlap_ratio=0.3, overlap_intensity=0.5,
                                               randomize_ratio=0.4, max_shift_ratio=0.3,
                                               resize_ratio=0.3, size_change_range=(0.4, 2.5),
                                               font_change_ratio=0.3, font_size_range=(-6, 8),
                                               preserve_title=True):
        """
        每页幻灯片应用多种扰动的组合（增强版）
        
        Args:
            delete_ratio: 删除元素比例
            overlap_ratio: 创建重叠的幻灯片比例
            overlap_intensity: 重叠强度
            randomize_ratio: 随机移动的元素比例
            max_shift_ratio: 最大位移比例
            resize_ratio: 调整文本框大小的比例
            size_change_range: 文本框大小变化范围（倍数）
            font_change_ratio: 改变字体的文本框比例
            font_size_range: 字体大小改变范围
            preserve_title: 是否保留标题
        
        Returns:
            dict: 包含各项扰动统计信息的字典
        """
        stats = {}
        
        print("开始应用PPT扰动（增强版：多种扰动组合）...")
        
        # 1. 删除元素
        print(f"1. 删除随机元素 (删除比例: {delete_ratio})...")
        deleted = self.delete_random_elements(delete_ratio, preserve_title)
        stats['deleted_elements'] = deleted
        print(f"   已删除 {deleted} 个元素")
        
        # 2. 创建元素重叠
        print(f"2. 创建元素重叠 (幻灯片比例: {overlap_ratio}, 强度: {overlap_intensity})...")
        overlapped = self.create_element_overlaps(overlap_ratio, overlap_intensity)
        stats['overlapped_slides'] = overlapped
        print(f"   已在 {overlapped} 个幻灯片中创建元素重叠")
        
        # 3. 随机移动元素
        print(f"3. 随机移动元素 (元素比例: {randomize_ratio}, 最大位移: {max_shift_ratio})...")
        moved = self.randomize_element_positions(randomize_ratio, max_shift_ratio)
        stats['moved_elements'] = moved
        print(f"   已随机移动 {moved} 个元素")
        
        # 4. 调整文本框大小
        print(f"4. 调整文本框大小 (比例: {resize_ratio}, 范围: {size_change_range})...")
        resized = self.resize_text_boxes(resize_ratio, size_change_range)
        stats['resized_textboxes'] = resized
        print(f"   已调整 {resized} 个文本框大小")
        
        # 5. 改变字体大小
        print(f"5. 改变字体大小 (比例: {font_change_ratio}, 范围: {font_size_range})...")
        font_changed = self.change_font_sizes(font_change_ratio, font_size_range)
        stats['font_changed'] = font_changed
        print(f"   已改变 {font_changed} 个文本框的字体大小")
        
        print("扰动完成！")
        return stats
    
    def apply_all_perturbations(self, delete_ratio=0.2, layout_change_ratio=0.3, 
                                font_change_ratio=0.3, size_change_range=(-4, 4),
                                preserve_title=True):
        """
        应用所有扰动（兼容旧版本的接口）
        
        Args:
            delete_ratio: 删除元素比例
            layout_change_ratio: 改变布局的幻灯片比例
            font_change_ratio: 改变字体的文本框比例
            size_change_range: 字体大小改变范围
            preserve_title: 是否保留标题
        
        Returns:
            dict: 包含各项扰动统计信息的字典
        """
        # 使用新的增强版方法
        return self.apply_combined_perturbations_per_slide(
            delete_ratio=delete_ratio,
            overlap_ratio=layout_change_ratio,
            overlap_intensity=0.5,
            randomize_ratio=layout_change_ratio * 1.3,
            max_shift_ratio=0.3,
            resize_ratio=font_change_ratio,
            size_change_range=(0.4, 2.5),
            font_change_ratio=font_change_ratio,
            font_size_range=size_change_range,
            preserve_title=preserve_title
        )
    
    def save(self, output_path):
        """
        保存扰动后的PPT
        
        Args:
            output_path: 输出文件路径
        """
        self.prs.save(output_path)
        print(f"已保存到: {output_path}")
